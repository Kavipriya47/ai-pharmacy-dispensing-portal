import sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx not found. Installing...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI-Assisted Pharmacy Medication Dispensing System"
subtitle.text = "An Adaptive Compliance Engine (ADCE) for Medication Safety, FEFO Control, and Regulatory Auditing"

# Slide 2: Problem Statement & Proposed Solution
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Problem Statement & Proposed Solution"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Traditional Manual Dispensing vs Proposed ADCE System:"
p = tf.add_paragraph()
p.text = "- Human Error Risk -> Active Safety Gate (Real-time rule validation)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Silent Safety Violations -> Instant Safety Block (Rejects expired/recalled stock)"
p.level = 1
p = tf.add_paragraph()
p.text = "- FEFO Bypasses -> Backend FEFO Engine (Auto-selects earliest-expiring batch)"
p.level = 1
p = tf.add_paragraph()
p.text = "- Paper-Based Tracing -> Instant Patient Tracing (Rapid contact tracing)"
p.level = 1

# Slide 3: System Architecture & Technology Stack
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "System Architecture & Technology Stack"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Three-Tier Architecture:"
p = tf.add_paragraph()
p.text = "Client Layer (Frontend)"
p.level = 1
p = tf.add_paragraph()
p.text = "React 19, TypeScript 6, Vite 8, Material-UI v9, TanStack Query 5"
p.level = 2
p = tf.add_paragraph()
p.text = "Service Layer (Backend)"
p.level = 1
p = tf.add_paragraph()
p.text = "Spring Boot 3.2 (Java 17), Spring Security, ADCE Safety Engine"
p.level = 2
p = tf.add_paragraph()
p.text = "Data Persistence & Compliance"
p.level = 1
p = tf.add_paragraph()
p.text = "Spring Data JPA, MySQL/H2 RDBMS"
p.level = 2

# Slide 4: Database ER Model & Schema Design
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Database ER Model & Schema Design"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Core Entities & Relationships:"
p = tf.add_paragraph()
p.text = "Supplier (1) --- (*) Medicine"
p.level = 1
p = tf.add_paragraph()
p.text = "Medicine (1) --- (1) Inventory"
p.level = 1
p = tf.add_paragraph()
p.text = "Inventory (1) --- (*) MedicineBatch"
p.level = 1
p = tf.add_paragraph()
p.text = "MedicineBatch (1) --- (*) DispensationRecord"
p.level = 1
p = tf.add_paragraph()
p.text = "AuditEvent: Independent log recording event type, user, timestamp, IP"
p.level = 1

# Slide 5: The ADCE Safety Engine (7 Business Rules)
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "The ADCE Safety Engine (7 Business Rules)"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Every dispense attempt passes through 7 real-time rules:"
p = tf.add_paragraph()
p.text = "1. Quantity strictly positive"
p.level = 1
p = tf.add_paragraph()
p.text = "2. Rejection of RECALLED stock"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Rejection of QUARANTINED stock"
p.level = 1
p = tf.add_paragraph()
p.text = "4. Rejection of non-ACTIVE status"
p.level = 1
p = tf.add_paragraph()
p.text = "5. Rejection of expired stock"
p.level = 1
p = tf.add_paragraph()
p.text = "6. Quantity stock availability check"
p.level = 1
p = tf.add_paragraph()
p.text = "7. Prescription enforcement (if required)"
p.level = 1

# Slide 6: FEFO Logic & Manual Override Workflow
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "FEFO Logic & Manual Override Workflow"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Default FEFO Mode:"
p = tf.add_paragraph()
p.text = "Backend auto-selects earliest-expiring active batch with sufficient stock"
p.level = 1
p = tf.add_paragraph()
p.text = "Frontend explicitly refrains from client-side prediction"
p.level = 1
p = tf.add_paragraph()
p.text = "Manual Batch Override Mode:"
p.level = 0
p = tf.add_paragraph()
p.text = "Pharmacist explicitly selects a specific batch"
p.level = 1
p = tf.add_paragraph()
p.text = "overrideReason becomes mandatory"
p.level = 1
p = tf.add_paragraph()
p.text = "System logs a FEFO_OVERRIDE audit event"
p.level = 1

# Slide 7: Batch Recall & Patient Tracing
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Batch Recall & Patient Tracing"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Workflow for critical safety recalls:"
p = tf.add_paragraph()
p.text = "1. Initiate Recall: Admin calls API endpoint for specific batch"
p.level = 1
p = tf.add_paragraph()
p.text = "2. Stock Adjustment: Status becomes RECALLED; remaining stock deducted"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Audit & Alert: Global CRITICAL severity notification; RECALL_INITIATED audit event"
p.level = 1
p = tf.add_paragraph()
p.text = "4. Patient Contact Tracing: Admin retrieves distinct list of affected patients"
p.level = 1

# Slide 8: Security Architecture
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Security Architecture & Single-Flight JWT Interceptor"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Dual-Layer Security Implementation:"
p = tf.add_paragraph()
p.text = "In-Memory Token Security: Access tokens held in JS memory (prevents XSS)"
p.level = 1
p = tf.add_paragraph()
p.text = "Stateless Authorization: Spring Security @PreAuthorize role guards"
p.level = 1
p = tf.add_paragraph()
p.text = "Single-Flight 401 Interceptor:"
p.level = 0
p = tf.add_paragraph()
p.text = "Intercepts 401 Unauthorized errors"
p.level = 1
p = tf.add_paragraph()
p.text = "Fires exactly 1 /auth/refresh request, queues parallel requests"
p.level = 1
p = tf.add_paragraph()
p.text = "Updates tokens and retries all queued requests seamlessly"
p.level = 1

# Slide 9: Reporting, Analytics & File Exports
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Reporting, Analytics & File Exports"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Real-Time Operations Analytics:"
p = tf.add_paragraph()
p.text = "Dispensing Summary & Inventory Dashboard KPIs"
p.level = 1
p = tf.add_paragraph()
p.text = "Role-Aware Security:"
p.level = 0
p = tf.add_paragraph()
p.text = "Restricted to ADMIN and AUDITOR roles"
p.level = 1
p = tf.add_paragraph()
p.text = "Native File Exports:"
p.level = 0
p = tf.add_paragraph()
p.text = "Export dispensing history to Excel (.xlsx) and PDF (.pdf)"
p.level = 1
p = tf.add_paragraph()
p.text = "Exported as raw binary blobs via ephemeral download links"
p.level = 1

# Slide 10: System Verification Results
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "System Verification Results"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Formal 43-Point System Hardening Matrix Passed:"
p = tf.add_paragraph()
p.text = "Phase H1: End-to-End Operational Workflow (11/11 Passed)"
p.level = 1
p = tf.add_paragraph()
p.text = "Phase H2: Role-Based Access Control & Security (15/15 Passed)"
p.level = 1
p = tf.add_paragraph()
p.text = "Phase H3: ADCE Safety Rules & JWT Refresh Resilience (10/10 Passed)"
p.level = 1
p = tf.add_paragraph()
p.text = "Phase H4: Production Readiness & Build Quality (7/7 Passed)"
p.level = 1
p = tf.add_paragraph()
p.text = "28/28 Backend Integration Tests Passed"
p.level = 1
p = tf.add_paragraph()
p.text = "0 Frontend Compilation/Build Errors"
p.level = 1

# Slide 11: Demonstration Roadmap
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "7-Step Live Demonstration Roadmap"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "System Walkthrough:"
p = tf.add_paragraph()
p.text = "1. Admin Login & Dashboard Overview"
p.level = 1
p = tf.add_paragraph()
p.text = "2. Catalog Setup (Supplier & Medicine)"
p.level = 1
p = tf.add_paragraph()
p.text = "3. Stock Intake (Batches with different expiries)"
p.level = 1
p = tf.add_paragraph()
p.text = "4. FEFO Dispensing & Rx Rule Enforcement"
p.level = 1
p = tf.add_paragraph()
p.text = "5. FEFO Manual Override (with audit event logging)"
p.level = 1
p = tf.add_paragraph()
p.text = "6. Batch Recall & Affected Patient Tracing"
p.level = 1
p = tf.add_paragraph()
p.text = "7. Reporting & Audit Export (Excel)"
p.level = 1

# Slide 12: Conclusion & Key Takeaways
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conclusion & Key Takeaways"
tf = slide.shapes.placeholders[1].text_frame
tf.text = "Summary:"
p = tf.add_paragraph()
p.text = "Active Safety Gate replaces passive inventory logs"
p.level = 1
p = tf.add_paragraph()
p.text = "FEFO Automation reduces inventory waste"
p.level = 1
p = tf.add_paragraph()
p.text = "Immutable audit logs ensure full traceability and regulatory compliance"
p.level = 1
p = tf.add_paragraph()
p.text = "Production-grade, secure architecture verified through rigorous E2E testing"
p.level = 1

prs.save('docs/AI_Assisted_Pharmacy_Defense_Presentation.pptx')
print("Presentation generated successfully at docs/AI_Assisted_Pharmacy_Defense_Presentation.pptx")
